"""Rebuild the Cheniere Sabine Pass deck in the exact visual format of the
Driftwood LNG deck (same theme, fonts, colors, table/card/badge layouts),
with added slides for Cheniere-specific economics (price architecture,
UFC/MSC, embedded Xy, cancellation/suspension mechanics) that have no
Driftwood analog.

Tables are built as plain shape grids (not native <a:tbl>), matching the
Driftwood deck's own approach and avoiding the corruption that makes
native tables unreadable by real PowerPoint (see export_slides.py).

Usage:
    python scripts/build_sabine_deck.py
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

ROOT = Path(__file__).parent.parent
DRIFTWOOD_PATH = ROOT / "driftwood_lng_spa_economic_terms_comparison_formula_signature_updated.pptx"
OUT_PATH = ROOT / "pptx" / "cheniere_sabine_pass_spa_standalone_economic_terms.pptx"

EMU_IN = 914400
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

NAVY = RGBColor(0x17, 0x32, 0x4D)
SLATE = RGBColor(0x47, 0x55, 0x69)
BODY = RGBColor(0x1E, 0x29, 0x3B)
BODY2 = RGBColor(0x1B, 0x2A, 0x3A)
MUTED = RGBColor(0x8A, 0x94, 0xA3)
BORDER = RGBColor(0xD7, 0xDE, 0xE8)
ROW_ALT = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2B, 0x5F, 0x8E)

RAG_BG = {
    "RED": RGBColor(0xF8, 0xD8, 0xD5),
    "AMBER": RGBColor(0xFF, 0xF0, 0xD5),
    "GREEN": RGBColor(0xDD, 0xEF, 0xE2),
    "N/A": RGBColor(0xE6, 0xE6, 0xE6),
    "NO RAG": RGBColor(0xE6, 0xE6, 0xE6),
}
RAG_TEXT = {
    "RED": RGBColor(0xB3, 0x26, 0x1E),
    "AMBER": RGBColor(0xB3, 0x6B, 0x00),
    "GREEN": RGBColor(0x2E, 0x7D, 0x32),
    "N/A": RGBColor(0x33, 0x33, 0x33),
    "NO RAG": RGBColor(0x1B, 0x2A, 0x3A),
}

FOOTER_TEXT = "Cheniere Sabine Pass SPA economic-terms comparison"


def in_(v):
    return Emu(int(v * EMU_IN))


def add_textbox(slide, left, top, width, height, text, size, color, bold=False, font="Calibri",
                 align=PP_ALIGN.LEFT, anchor=None, wrap=True, line_spacing=None):
    box = slide.shapes.add_textbox(in_(left), in_(top), in_(width), in_(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill=None, line_color=None, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, in_(left), in_(top), in_(width), in_(height))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_header(slide, title, subtitle):
    add_textbox(slide, 0.55, 0.28, 10.8, 0.42, title, 20, NAVY, bold=True, font="Aptos Display")
    add_textbox(slide, 0.58, 0.76, 11.7, 0.3, subtitle, 9.5, SLATE)
    line = add_rect(slide, 0.55, 1.12, 12.2, 0.0, fill=None, line_color=BORDER, line_w=1)


def add_footer(slide, page_num):
    add_textbox(slide, 0.55, 7.14, 6.8, 0.18, FOOTER_TEXT, 7.5, MUTED)
    add_textbox(slide, 12.35, 7.13, 0.4, 0.2, str(page_num), 8.5, MUTED)


def new_content_slide(prs, blank_layout, title, subtitle, page_num):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, title, subtitle)
    add_footer(slide, page_num)
    return slide


def add_title_slide(prs, blank_layout, title, subtitle, tagline, purpose_label, purpose_text,
                     legend_items, page_num, arc_shapes):
    slide = prs.slides.add_slide(blank_layout)
    bg = add_rect(slide, 0, 0, 13.333, 7.5, fill=WHITE, line_color=RGBColor(0x33, 0x33, 0x33), line_w=1)
    add_textbox(slide, 0.72, 0.82, 7.1, 0.62, title, 31, NAVY, bold=True, font="Aptos Display")
    add_textbox(slide, 0.75, 1.55, 7.9, 0.36, subtitle, 15, SLATE)
    add_textbox(slide, 0.75, 2.12, 7.2, 0.26, tagline, 10.5, BODY)
    add_textbox(slide, 0.75, 3.12, 1.3, 0.25, purpose_label, 11, BLUE, bold=True)
    add_textbox(slide, 0.75, 3.47, 6.1, 1.4, purpose_text, 13.2, BODY)
    for shp_xml in arc_shapes:
        slide.shapes._spTree.append(copy.deepcopy(shp_xml))
    colors = [BLUE, RGBColor(0xB3, 0x6B, 0x00), RGBColor(0xB3, 0x26, 0x1E)]
    xs = [8.4, 9.82, 11.55]
    widths = [1.3, 1.65, 1.3]
    for x, w, label, color in zip(xs, widths, legend_items, colors):
        add_textbox(slide, x, 5.4, w, 0.24, label, 9, color, bold=True)
    add_footer(slide, page_num)
    return slide


def add_three_cards(slide, cards, top=1.45):
    xs = [0.65, 4.85, 9.05]
    widths = [3.75, 3.75, 3.65]
    for x, w, (label, color, text) in zip(xs, widths, cards):
        add_textbox(slide, x, top, w, 0.28, label, 11.5, color, bold=True)
        add_textbox(slide, x, top + 0.36, w, 1.4, text, 10.5, BODY)


def add_rag_legend_table(slide, rows, top=3.65, row_h=0.55):
    col_x = [0.8, 2.0, 4.15]
    col_w = [1.2, 2.15, 8.2]
    headers = ["RAG", "Commercial meaning", "Key items"]
    for x, w, h in zip(col_x, col_w, headers):
        add_rect(slide, x, top, w, row_h, fill=NAVY)
        add_textbox(slide, x + 0.08, top, w - 0.16, row_h, h, 10, WHITE, bold=True,
                    anchor=3)
    y = top + row_h
    for i, (rag, meaning, items) in enumerate(rows):
        alt = i % 2 == 1
        base_fill = ROW_ALT if alt else WHITE
        rag_bg = RAG_BG.get(rag, base_fill)
        rag_fg = RAG_TEXT.get(rag, BODY)
        add_rect(slide, col_x[0], y, col_w[0], row_h, fill=rag_bg, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[0] + 0.08, y, col_w[0] - 0.16, row_h, rag, 10, rag_fg, bold=True, anchor=3)
        add_rect(slide, col_x[1], y, col_w[1], row_h, fill=base_fill, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[1] + 0.08, y, col_w[1] - 0.16, row_h, meaning, 10, BODY, anchor=3)
        add_rect(slide, col_x[2], y, col_w[2], row_h, fill=base_fill, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[2] + 0.08, y, col_w[2] - 0.16, row_h, items, 8.5, BODY2, anchor=3)
        y += row_h


def add_two_col_table(slide, headers, rows, top=1.42, row_h=0.55, col1_w=2.45, font_size=10.2):
    col_x = [0.72, 3.17]
    col_w = [col1_w, 12.2 - 0.72 - col1_w - 0.15]
    for x, w, h in zip(col_x, col_w, headers):
        add_rect(slide, x, top, w, row_h, fill=NAVY)
        add_textbox(slide, x + 0.08, top, w - 0.16, row_h, h, 10, WHITE, bold=True, anchor=3)
    y = top + row_h
    for i, (label, desc) in enumerate(rows):
        alt = i % 2 == 1
        fill = ROW_ALT if alt else WHITE
        add_rect(slide, col_x[0], y, col_w[0], row_h, fill=fill, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[0] + 0.08, y, col_w[0] - 0.16, row_h, label, font_size, BODY, anchor=3)
        add_rect(slide, col_x[1], y, col_w[1], row_h, fill=fill, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[1] + 0.08, y, col_w[1] - 0.16, row_h, desc, font_size, BODY, anchor=3)
        y += row_h
    return y


def add_rag_table(slide, headers, rows, top=1.45, row_h=0.56, first_col_w=1.35, rag_col_w=0.9,
                   font_size=9.2, header_font_size=9.2):
    n_data_cols = len(headers) - 2
    total_w = 12.2
    remaining = total_w - first_col_w - rag_col_w - 0.15 * (len(headers) - 1)
    data_col_w = remaining / n_data_cols
    col_x = [0.46]
    col_w = [first_col_w]
    for i in range(n_data_cols):
        col_x.append(col_x[-1] + col_w[-1] + 0.15 / len(headers) * 0 + 0.0)
    # simpler: lay out contiguous columns with a small gap
    gap = 0.0
    xs = [0.46]
    widths = [first_col_w] + [data_col_w] * n_data_cols + [rag_col_w]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)

    for x, w, h in zip(xs, widths, headers):
        add_rect(slide, x, top, w, row_h, fill=NAVY)
        add_textbox(slide, x + 0.06, top, w - 0.12, row_h, h, header_font_size, WHITE, bold=True, anchor=3)

    y = top + row_h
    for i, row in enumerate(rows):
        alt = i % 2 == 1
        fill = ROW_ALT if alt else WHITE
        for j, (x, w) in enumerate(zip(xs, widths)):
            val = row[j]
            is_rag_col = j == len(widths) - 1
            if is_rag_col:
                rag_bg = RAG_BG.get(val, fill)
                rag_fg = RAG_TEXT.get(val, BODY)
                add_rect(slide, x, y, w, row_h, fill=rag_bg, line_color=BORDER, line_w=0.5)
                add_textbox(slide, x + 0.06, y, w - 0.12, row_h, val, font_size, rag_fg, bold=True, anchor=3)
            else:
                add_rect(slide, x, y, w, row_h, fill=fill, line_color=BORDER, line_w=0.5)
                add_textbox(slide, x + 0.06, y, w - 0.12, row_h, str(val), font_size, BODY, anchor=3)
        y += row_h
    return y


def add_badge_rows(slide, sections, top=1.35):
    y = top
    for section_title, items in sections:
        if section_title:
            add_textbox(slide, 0.75, y, 11.6, 0.24, section_title, 10.8, NAVY, bold=True)
            y += 0.36
        for rag, label, desc in items:
            bg = RAG_BG.get(rag, RGBColor(0xE6, 0xE6, 0xE6))
            fg = RAG_TEXT.get(rag, BODY)
            add_rect(slide, 0.75, y + 0.02, 0.62, 0.2, fill=bg)
            add_textbox(slide, 0.75, y + 0.02, 0.62, 0.2, rag, 7.5, fg, bold=True, align=PP_ALIGN.CENTER, anchor=3)
            add_textbox(slide, 1.49, y, 2.55, 0.25, label, 8.6, BODY, bold=True)
            add_textbox(slide, 4.1, y, 8.25, 0.28, desc, 8.4, SLATE)
            y += 0.37
    return y


def add_callout(slide, top, text, fill=RGBColor(0xEE, 0xF4, 0xF8), height=0.62):
    add_rect(slide, 0.74, top, 11.6, height, fill=fill)
    add_textbox(slide, 0.86, top, 11.4, height, text, 11.2, BODY2, anchor=3)


def add_formula_boxes(slide, boxes, top=1.25, height=2.1):
    xs = [0.55, 6.9]
    widths = [6.15, 5.9]
    for x, w, (fill, label, body) in zip(xs, widths, boxes):
        add_rect(slide, x, top, w, height, fill=fill, line_color=RGBColor(0xCB, 0xD5, 0xE1), line_w=0.75)
        add_textbox(slide, x + 0.12, top + 0.08, w - 0.24, 0.22, label, 9, BODY, bold=True)
        add_textbox(slide, x + 0.12, top + 0.36, w - 0.24, height - 0.44, body, 7.6, RGBColor(0x33, 0x41, 0x55))


def get_arc_shapes():
    prs = Presentation(DRIFTWOOD_PATH)
    slide1 = list(prs.slides)[0]
    arcs = []
    for shape in slide1.shapes:
        if shape.name in ("Shape 6", "Shape 7", "Shape 8"):
            arcs.append(shape._element)
    return arcs


def build():
    arc_shapes = get_arc_shapes()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    page = 1

    # ---- Slide 1: Title ----
    add_title_slide(
        prs, blank,
        "Cheniere Sabine Pass LNG SPAs",
        "Standalone economic-term comparison and valuation fields",
        "BG  •  Gas Natural Fenosa  •  GAIL  •  Total",
        "Purpose",
        "Identify how fixed liquefaction economics are recovered — separate Unit Fixed "
        "Charge / Monthly Sales Charge, or embedded Xy inside CSP — and whether those "
        "charges survive cancellation, suspension or non-lifting. This is an economic "
        "comparison, not a legal blackline.",
        ["Visible terms", "Redactions / unknowns", "Modelling impact"],
        page, arc_shapes,
    )
    page += 1

    # ---- Slide 2: Executive read-across ----
    slide = new_content_slide(
        prs, blank, "Executive read-across",
        "The contracts use a common Cheniere Sabine Pass SPA skeleton but differ sharply "
        "in how fixed liquefaction economics are recovered.",
        page,
    )
    add_three_cards(slide, [
        ("Fact", RGBColor(0x2E, 0x7D, 0x32),
         "All four contracts are FOB LNG SPAs from the Sabine Pass Facility with Henry "
         "Hub-linked commodity pricing (CSP = 1.15 × HH), common 20-year base term "
         "and broadly common delivery, title/risk and credit-support architecture."),
        ("Inference", BLUE,
         "Model a dedicated price-architecture and capacity-payment section: the fixed/"
         "liquefaction component is recovered differently — separate UFC/MSC (BG, GNF) "
         "versus an Xy component embedded in CSP (GAIL, Total)."),
        ("Unknown", RGBColor(0xB3, 0x6B, 0x00),
         "CPI baselines and escalation weights are visible, but buyer-specific credit "
         "thresholds, LC mechanics and full DoP/cover-damages formulas need clause-by-"
         "clause review before they are treated as equivalent."),
    ])
    add_rag_legend_table(slide, [
        ("GREEN", "No material difference detected",
         "Base HH-linked commodity index, FOB delivery, title/risk, buyer shipping"),
        ("RED", "Visible economic difference",
         "Price architecture (UFC/MSC vs embedded Xy), CPI escalation weight, "
         "cancellation/suspension mechanics, delivery shape"),
        ("AMBER", "Framework similar, but caveats",
         "Extension conditions, credit support, maintenance/inspection caps, off-spec remedies"),
    ], top=3.55, row_h=0.5)
    page += 1

    # ---- Slide 3: Contract universe ----
    slide = new_content_slide(
        prs, blank, "Contract universe",
        "ACQ is shown but not RAG-scored; size is buyer-specific commercial volume, not a drafting difference.",
        page,
    )
    y = add_rag_table(
        slide,
        ["Contract", "Seller", "Buyer", "Signature / effective date", "Visible ACQ", "RAG"],
        [
            ["BG", "Sabine Pass Liquefaction", "BG Gulf Coast LNG", "25 Oct 2011", "182.5m MMBtu/yr", "NO RAG"],
            ["Gas Natural Fenosa", "Sabine Pass Liquefaction", "Gas Natural Aprovisionamientos",
             "21 Nov 2011", "182.5m MMBtu/yr", "NO RAG"],
            ["GAIL", "Sabine Pass Liquefaction", "GAIL (India) Limited", "11 Dec 2011",
             "182.5m MMBtu/yr", "NO RAG"],
            ["Total", "Sabine Pass Liquefaction", "Total Gas & Power North America", "14 Dec 2012",
             "104.75m MMBtu/yr = 91.25m AT + 13.5m ST", "NO RAG"],
        ],
        first_col_w=1.9, font_size=9.5,
    )
    add_callout(slide, y + 0.25,
                "Observation: Total is structurally different because ACQ is the sum of an Annual "
                "Tranche and a Seasonal Tranche. The Seasonal Tranche is 13.5m MMBtu, approximately "
                "12.9% of Total ACQ.  AT = Annual Tranche.  ST = Seasonal Tranche.",
                height=0.62)
    page += 1

    # ---- Slide 4 (NEW/Cheniere-specific): Price architecture & capacity-payment fields ----
    slide = new_content_slide(
        prs, blank, "Price architecture and capacity-payment fields",
        "Cheniere-specific fields required to analyse how each agreement separates or embeds liquefaction economics.",
        page,
    )
    y = add_two_col_table(slide, ["Driver", "Why it is needed"], [
        ("Price architecture", "Distinguishes CSP-only from CSP + UFC + Monthly Sales Charge"),
        ("Unit Fixed Charge / UFC formula", "BG and Gas Natural have a separate CPI-linked fixed component"),
        ("Monthly Sales Charge", "Capacity reservation charge payable monthly, not cargo-by-cargo"),
        ("X0 / Xy fixed component", "GAIL and Total embed the fixed component in CSP"),
        ("CPI escalation weight", "GAIL uses 15%; Total uses 11.5%; BG/GNF use the UFC formula"),
        ("Charge survives cancellation / suspension?", "Core economic distinction for cancelled or suspended cargoes"),
        ("Early / bridging / seasonal volumes", "BG has Early Cargoes; GAIL has Bridging Volumes; Total has a seasonal tranche"),
        ("Damages reference price", "BG/GNF refer to CSP minus UFC; GAIL/Total use CSP-based mechanics"),
    ], row_h=0.5, font_size=9.6)
    add_callout(slide, y + 0.18,
                "Retain standard SPA fields for CPs, term, delivery point, title/risk, ADP, quality, FM, "
                "credit and taxes. Add the fields above before comparing Cheniere agreements or other "
                "US liquefaction-style LNG SPAs.", height=0.62)
    page += 1

    # ---- Slide 5: Visible price formulas ----
    slide = new_content_slide(
        prs, blank, "Visible price formulas",
        "All four use 1.15 × HH, but the fixed/liquefaction component is recovered differently.",
        page,
    )
    y = add_rag_table(
        slide,
        ["Contract", "Visible formula", "Fixed / capacity component", "RAG"],
        [
            ["BG", "CSP = 1.15 × HH", "UFCy = 1.9125 + (0.3375 × CPIy / CPI0); MSC = UFC × Q / 12", "RED"],
            ["Gas Natural Fenosa", "CSP = 1.15 × HH", "UFCy = 2.1525 + (0.3375 × CPIy / CPI0); MSC = UFC × Q / 12", "RED"],
            ["GAIL", "CSP = (1.15 × HH) + Xy", "Xy = (0.85 + 0.15 × CPI(y-1) / CPI0) × X0; X0 = US$3.00/MMBtu", "RED"],
            ["Total", "CSP = (1.15 × HH) + Xy", "Xy = (0.885 + 0.115 × CPI(y-1) / CPI0) × X0; X0 = US$3.00/MMBtu", "RED"],
        ],
        first_col_w=1.9, rag_col_w=0.85, font_size=9.3,
    )
    add_callout(slide, y + 0.25,
                "RAG is RED because formula architecture differs, not because each contract has a "
                "different HH multiplier. The HH multiplier is common at 1.15. For modelling: separate "
                "variable gas index exposure from fixed-capacity economics before calculating margin, "
                "cancellation value or DoP/cover damages.", height=0.75)
    page += 1

    # ---- Slide 6 (NEW appendix): Base formula text as filed ----
    slide = new_content_slide(
        prs, blank, "Appendix A: base formula text as filed",
        "Visible and redacted formula wording. CPI baseline constants remain redacted as [***] where applicable.",
        page,
    )
    add_formula_boxes(slide, [
        (RGBColor(0xDB, 0xEA, 0xFE), "BG unit fixed charge",
         "UFCy = 1.9125 + (0.3375 × CPIy / CPI0)\nMSC = UFC × Q / 12\n\n"
         "CSP = 1.15 × HH. The UFC is a separate, CPI-escalated capacity charge "
         "invoiced monthly against scheduled quantity."),
        (RGBColor(0xDB, 0xEA, 0xFE), "Gas Natural Fenosa unit fixed charge",
         "UFCy = 2.1525 + (0.3375 × CPIy / CPI0)\nMSC = UFC × Q / 12\n\n"
         "Same architecture as BG; base constant (2.1525) differs from BG's 1.9125."),
    ], top=1.25, height=1.95)
    add_formula_boxes(slide, [
        (RGBColor(0xFE, 0xE2, 0xE2), "GAIL embedded Xy",
         "CSP = (1.15 × HH) + Xy\nXy = (0.85 + 0.15 × CPI(y-1) / CPI0) × X0\nX0 = US$3.00/MMBtu\n\n"
         "The fixed component is embedded inside the cargo price itself, not invoiced separately."),
        (RGBColor(0xFE, 0xE2, 0xE2), "Total embedded Xy",
         "CSP = (1.15 × HH) + Xy\nXy = (0.885 + 0.115 × CPI(y-1) / CPI0) × X0\nX0 = US$3.00/MMBtu\n\n"
         "Same embedded-Xy architecture as GAIL; CPI weight differs (11.5% vs 15%)."),
    ], top=3.55, height=1.95)
    add_callout(slide, 5.75,
                "Fact: the HH multiplier (1.15) and X0 (US$3.00/MMBtu) are common. Unknown: whether "
                "CPI base-year conventions are identical across contracts — do not infer equality "
                "of embedded value from formula shape alone.", height=0.7)
    page += 1

    # ---- Slide 7: Liquefaction-margin recovery ----
    slide = new_content_slide(
        prs, blank, "Liquefaction-margin recovery",
        "This is the main economic distinction across the Cheniere agreements.",
        page,
    )
    add_textbox(slide, 0.75, 1.5, 5.6, 0.28, "BG / Gas Natural Fenosa", 12, NAVY, bold=True)
    add_textbox(slide, 6.9, 1.5, 5.6, 0.28, "GAIL / Total", 12, NAVY, bold=True)
    add_rect(slide, 0.75, 1.95, 5.6, 1.7, fill=RGBColor(0xF5, 0xF7, 0xFA))
    add_textbox(slide, 0.95, 2.1, 5.2, 1.4,
                "Cargo price\nCSP = 1.15 × HH\n\nCapacity charge\nUFC × Q / 12",
                12, BODY, line_spacing=1.3)
    add_rect(slide, 6.9, 1.95, 5.6, 1.7, fill=RGBColor(0xF5, 0xF7, 0xFA))
    add_textbox(slide, 7.1, 2.1, 5.2, 1.4,
                "Cargo price\nCSP = 1.15 × HH + Xy\n\nFixed component\ninside CSP",
                12, BODY, line_spacing=1.3)
    add_callout(slide, 3.95,
                "Economic consequence: BG/GNF capacity payments can continue irrespective of the "
                "cargo actually taken. GAIL/Total embed the fixed element in the cargo price, so "
                "cancellation/suspension mechanics must be read differently.", height=0.75)
    add_callout(slide, 4.95,
                "Standalone section required: price architecture and capacity-payment mechanics.",
                fill=RGBColor(0xF5, 0xF7, 0xFA), height=0.5)
    page += 1

    # ---- Slide 8: Term and extension mechanics ----
    slide = new_content_slide(
        prs, blank, "Term and extension mechanics",
        "Base term is broadly common; extension mechanics have drafting differences and approvals constraints.",
        page,
    )
    add_rag_table(slide, ["Item", "BG", "Gas Natural Fenosa", "GAIL", "Total", "RAG"], [
        ["Base term", "20 years from DFCD", "20 years from DFCD", "20 years from DFCD", "20 years from DFCD", "GREEN"],
        ["Extension", "Buyer may extend up to 10 years for any portion of ACQ", "Similar maximum extension term",
         "Buyer may extend up to 10 years", "Buyer may extend up to 10 years", "AMBER"],
        ["Extension condition", "Minimum facility utilisation / approvals and export authorisations", "Similar",
         "Approvals; no same export-authorisation wording visible", "Approvals incl. LNG export licences", "AMBER"],
        ["Delayed DFCD termination", "Buyer right after delay; seller has corresponding mechanics", "Similar",
         "Similar, but first-window timing differs", "Similar; 455-day FM deferral cap visible", "AMBER"],
    ], first_col_w=1.7, font_size=8.4, row_h=0.72)
    page += 1

    # ---- Slide 9: Volume, shape and maintenance ----
    slide = new_content_slide(
        prs, blank, "Volume, shape and maintenance: percentage view",
        "ACQ itself is unscored; numerical caps are compared as percentage of ACQ where possible.",
        page,
    )
    add_rag_table(slide, ["Item", "BG", "Gas Natural Fenosa", "GAIL", "Total", "RAG"], [
        ["ACQ size", "182.5m MMBtu", "182.5m MMBtu", "182.5m MMBtu", "104.75m MMBtu", "NO RAG"],
        ["Delivery shape", "Even / ratable", "Even / ratable", "Even / ratable", "Annual + Seasonal Tranche", "RED"],
        ["Seasonal share", "n/a", "n/a", "n/a", "13.5m / 104.75m = 12.9% of ACQ", "RED"],
        ["Major maintenance cap", "6% per year; 30% initial term", "6% per year; 30% initial term",
         "6% per year; 30% initial term", "6% per year; 30% initial term", "GREEN"],
        ["Inspection cap", "8.3% per year; 16.6% initial term", "8.3% per year; 16.6% initial term",
         "8.3% per year; 16.6% initial term", "8.3% of annual/seasonal base; 17% initial term", "AMBER"],
    ], first_col_w=1.7, font_size=8.2, row_h=0.6)
    page += 1

    # ---- Slide 10: Cancellation, suspension and early-volume mechanics ----
    slide = new_content_slide(
        prs, blank, "Cancellation, suspension and early-volume mechanics",
        "These are economically central in Cheniere because fixed charges may survive non-lifting.",
        page,
    )
    y = add_rag_table(slide, ["Item", "BG", "Gas Natural Fenosa", "GAIL", "Total", "RAG"], [
        ["Buyer cancellation / suspension", "Cancellation right; Monthly Sales Charge continues",
         "Cancellation right; Monthly Sales Charge continues", "No equivalent seen in reviewed snippets",
         "Suspension right with Suspension Fee", "RED"],
        ["Early / bridging volumes", "Early Cargoes; pre-commercial cargo pricing seen", "No equivalent highlighted",
         "Bridging Volumes", "Early First Window Period", "RED"],
        ["Capacity payment if cargo cancelled", "Yes: MSC continues", "Yes: MSC continues", "n/a in same form",
         "Suspension fee mechanics", "RED"],
        ["FM affecting LNG tanker / discharge terminal", "Express clause", "Express clause",
         "Express clause; terminal wording narrower", "Express clause", "AMBER"],
    ], first_col_w=2.3, font_size=7.8, row_h=0.72)
    add_callout(slide, y + 0.15,
                "Do not value cancellation as simply “no cargo, no price”. In BG/GNF, capacity "
                "reservation payments are a separate cashflow stream. For Total, suspension requires a "
                "separate Suspension Fee analysis rather than a shortfall-only treatment.", height=0.62)
    page += 1

    # ---- Slide 11: Damages and remedy fields ----
    slide = new_content_slide(
        prs, blank, "Damages and remedy fields to add",
        "The relevant reference price differs where fixed charges are separated from CSP.",
        page,
    )
    y = add_rag_table(slide, ["Economic item", "BG", "Gas Natural Fenosa", "GAIL", "Total", "RAG"], [
        ["Buyer shortfall / cover damages", "CSP-based; cancellation interacts with MSC", "Same broad architecture",
         "CSP-based cover damages", "CSP-based plus suspension / idling-cost wording", "RED"],
        ["Seller DoP", "Formula refers to CSP minus UFC for replacement comparison", "Same broad concept",
         "CSP/Cover Damages mechanics", "CSP plus additional cost wording", "RED"],
        ["Monthly / capacity charge in remedies", "MSC continues / is separately invoiced",
         "MSC continues / separately invoiced", "Not separate", "Suspension Fee separate", "RED"],
        ["Off-spec LNG remedies", "Similar SPA architecture", "Similar", "Similar", "Similar", "AMBER"],
    ], first_col_w=2.4, font_size=7.6, row_h=0.72)
    add_callout(slide, y + 0.15,
                "New rows to add to the master comparison: damages reference price; capacity-charge "
                "survival; seller DoP adjustment for UFC; suspension-fee calculation; invoice type "
                "for capacity charges.", height=0.6)
    page += 1

    # ---- Slide 12: Delivery and logistics ----
    slide = new_content_slide(
        prs, blank, "Delivery and logistics",
        "Core FOB structure is common across the four agreements, but some infrastructure-specific schedules and exhibits differ.",
        page,
    )
    y = add_rag_table(slide, ["Item", "BG", "Gas Natural Fenosa", "GAIL", "Total", "RAG"], [
        ["Delivery point / incoterm", "FOB Sabine Pass Facility", "FOB Sabine Pass Facility",
         "FOB Sabine Pass Facility", "FOB Sabine Pass Facility", "GREEN"],
        ["Buyer shipping", "Buyer provides tankers", "Same", "Same", "Same", "GREEN"],
        ["Destination flexibility", "Buyer may market to any destination, subject to agreement",
         "Same broad wording", "Same broad wording", "Same broad wording", "GREEN"],
        ["Alternate source delivery", "Permitted with buyer consent and conditions",
         "Similar; Gulf Coast affiliate carve-out visible", "Similar", "Similar", "AMBER"],
        ["Infrastructure exhibits", "Terminal/cooperation/tug/transport exhibits visible",
         "Not same set highlighted", "Not same set highlighted", "Not same set highlighted", "AMBER"],
    ], first_col_w=2.1, font_size=8.0, row_h=0.6)
    add_callout(slide, y + 0.15,
                "Keep standard logistics rows, and add a Cheniere-specific field for linked "
                "infrastructure exhibits where they affect rights, costs or facility access.",
                height=0.5)
    page += 1

    # ---- Slide 13: Credit, financing and bankability ----
    slide = new_content_slide(
        prs, blank, "Credit, financing and bankability fields",
        "These need full clause-by-clause checks before they are treated as equivalent.",
        page,
    )
    y = add_rag_table(slide, ["Item", "BG", "Gas Natural Fenosa", "GAIL", "Total", "RAG"], [
        ["Buyer credit support", "Credit support clause present", "Present", "Present", "Present", "AMBER"],
        ["Parent / guaranty mechanics", "Buyer-specific", "Buyer-specific", "Buyer-specific", "Buyer-specific", "AMBER"],
        ["Lender / direct agreement mechanics", "Project-finance style provisions", "Similar", "Similar", "Similar", "AMBER"],
        ["Limitations on liability", "Caps/limits present", "Similar", "Similar", "Similar", "AMBER"],
        ["Representations / business practices", "Present", "Present", "Present", "Present", "AMBER"],
    ], first_col_w=2.2, font_size=8.2, row_h=0.6)
    add_callout(slide, y + 0.15,
                "Action: do not assume credit-support equivalence from the template. Buyer-specific "
                "parent support, LC triggers and credit ratings should be extracted separately.",
                height=0.5)
    page += 1

    # ---- Slide 14: First-pass RAG summary ----
    slide = new_content_slide(
        prs, blank, "First-pass RAG summary",
        "High-confidence economic differences are concentrated in pricing, capacity charges, volume shape and non-lifting rights.",
        page,
    )
    add_rag_table(slide, ["Economic item", "BG", "Gas Natural Fenosa", "GAIL", "Total", "RAG"], [
        ["Signature / effective date", "25 Oct 2011", "21 Nov 2011", "11 Dec 2011", "14 Dec 2012", "NO RAG"],
        ["Base commodity index", "HH", "HH", "HH", "HH", "GREEN"],
        ["Price architecture", "HH + UFC/MSC", "HH + UFC/MSC", "HH + Xy", "HH + Xy", "RED"],
        ["CPI escalation weight", "UFC formula", "UFC formula", "15% of X0", "11.5% of X0", "RED"],
        ["Delivery shape", "Ratable", "Ratable", "Ratable", "Annual + seasonal", "RED"],
        ["Buyer non-lifting right", "Cancellation", "Cancellation", "Not same", "Suspension", "RED"],
        ["Maintenance percentages", "6%/30%; 8.3%/16.6%", "Same", "Same", "Inspection base differs; 17% cumulative", "AMBER"],
        ["DoP / damages basis", "CSP-UFC interaction", "CSP-UFC interaction", "CSP", "CSP + suspension/idling elements", "RED"],
    ], first_col_w=2.1, font_size=7.4, row_h=0.6)
    page += 1

    # ---- Slide 15: Recommended valuation object model ----
    slide = new_content_slide(
        prs, blank, "Recommended valuation object model",
        "Separate cargo economics from capacity economics before computing value, hedge and risk.",
        page,
    )
    add_three_cards(slide, [
        ("1. Variable commodity leg", BLUE, "HH exposure: 1.15 × HH × loaded MMBtu."),
        ("2. Fixed liquefaction leg", RGBColor(0xB3, 0x6B, 0x00),
         "UFC/MSC or Xy component, CPI escalated."),
        ("3. Non-lifting logic", RGBColor(0xB3, 0x26, 0x1E),
         "Cancellation, suspension, MSC survival, suspension fee."),
    ], top=1.5)
    add_callout(slide, 3.15,
                "Inputs to store per contract: ACQ and tranche split; DFCD and term; extension volume; "
                "HH multiplier; UFC/X0/Xy formula; CPI base; monthly charge formula; cancellation/"
                "suspension right; damage reference price; FM/capacity-charge reduction; credit "
                "support; delivery profile.", fill=RGBColor(0xF5, 0xF7, 0xFA), height=1.1)
    add_callout(slide, 4.45,
                "Output metrics: gross margin by leg; fixed-capacity recovery; hedge exposure to HH and "
                "CPI; optionality/cancellation exposure; DoP/cover damages under stress cases; "
                "credit-adjusted PV.", fill=RGBColor(0xEE, 0xF4, 0xF8), height=0.85)
    page += 1

    # ---- Slide 16: Caveats and next steps ----
    slide = new_content_slide(
        prs, blank, "Caveats and next steps",
        "This is an economic comparison; it is not a legal blackline.",
        page,
    )
    add_badge_rows(slide, [
        (None, [
            ("AMBER", "Buyer credit-support thresholds", "Extract LC mechanics — buyer-specific credit economics may be material."),
            ("AMBER", "BG vs Gas Natural UFC/MSC blackline", "The UFC base differs visibly; remedy/cancellation interaction matters."),
            ("AMBER", "GAIL Bridging Volumes / Total suspension fee", "Read in full — potentially material transition-period economics."),
            ("AMBER", "Full DoP / cover damages formulas", "Reference price differs where UFC is separate."),
            ("AMBER", "Cheniere-specific template fields", "Add to the master LNG SPA comparison template; without them the "
             "table understates capacity-payment economics."),
        ]),
    ], top=1.5)
    add_callout(slide, 4.0,
                "Main structural addition: add “Price architecture and capacity-payment mechanics” "
                "before the pricing RAG table.", fill=RGBColor(0xF5, 0xF7, 0xFA), height=0.6)
    page += 1

    OUT_PATH.parent.mkdir(exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved {OUT_PATH} ({page - 1} slides)")


if __name__ == "__main__":
    build()
