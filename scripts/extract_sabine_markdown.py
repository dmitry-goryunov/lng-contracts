"""Extract the Cheniere Sabine Pass deck (built by build_sabine_deck.py) to a
readable Markdown file, one '## Slide N: Title' section per slide, matching
the format of the other converted decks in contracts/.

Usage:
    python scripts/extract_sabine_markdown.py
"""

from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).parent.parent
PPTX_PATH = ROOT / "pptx" / "cheniere_sabine_pass_spa_standalone_economic_terms.pptx"
OUT_PATH = ROOT / "contracts" / "cheniere_sabine_pass_spa_standalone_economic_terms.md"

FOOTER_TEXT = "Cheniere Sabine Pass SPA economic-terms comparison"


def slide_title(slide, index):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if text and text != FOOTER_TEXT and not text.isdigit():
            return text
    return f"Slide {index}"


def slide_body_lines(slide, title_text):
    seen_title = False
    lines = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text or text == FOOTER_TEXT or text.isdigit():
            continue
        if not seen_title and text == title_text:
            seen_title = True
            continue
        lines.append(text)
    return lines


def build():
    prs = Presentation(PPTX_PATH)
    slides = list(prs.slides)

    out = [
        "# cheniere sabine pass spa standalone economic terms",
        "",
        f"Source file: `{PPTX_PATH.name}`",
        "",
        "Converted from PowerPoint to Markdown. Slide order and visible text are preserved "
        "as a readable text version.",
        "",
        "---",
        "",
    ]

    for i, slide in enumerate(slides, start=1):
        title = slide_title(slide, i)
        out.append(f"## Slide {i}: {title}")
        out.append("")
        for line in slide_body_lines(slide, title):
            out.append(line)
            out.append("")
        out.append("---")
        out.append("")

    OUT_PATH.write_text("\n".join(out), encoding="utf-8")
    print(f"wrote {OUT_PATH} ({len(slides)} slides)")


if __name__ == "__main__":
    build()
