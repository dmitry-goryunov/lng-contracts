"""Build the Cheniere follow-on & GdF deck in the exact visual format of the
Driftwood / Sabine Pass decks (same theme, fonts, colors, table/card/badge
layouts). Covers the five contracts not in the first two decks:

  11  Sabine Pass - BG Gulf Coast, Amended & Restated (25 Jan 2012)
  13  Sabine Pass - Centrica (22 Mar 2013)
  12  Corpus Christi - Gas Natural Fenosa (2 Jun 2014)
  14  Corpus Christi - Woodside (30 Jun 2014)
  15  Cheniere Marketing - Gaz de France Master Ex-Ship (26 Apr 2007)

Tables are plain shape grids (not native <a:tbl>), matching the other decks
and avoiding the corruption that makes native tables unreadable by real
PowerPoint (see export_slides.py).

Title-slide arc decorations are copied from the Driftwood deck when that file
is present next to this repo (exact visual match); otherwise approximated
with drawn arc shapes so the script also runs without the binary source deck.

All values are taken from the full contract texts in contracts/. Redacted or
unextracted values are never guessed.

Usage:
    python scripts/build_followon_deck.py
"""

import copy
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

ROOT = Path(__file__).parent.parent
DRIFTWOOD_PATH = ROOT / "driftwood_lng_spa_economic_terms_comparison_formula_signature_updated.pptx"
OUT_PATH = ROOT / "pptx" / "cheniere_followon_gdf_spa_economic_terms.pptx"

EMU_IN = 914400
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

NAVY = RGBColor(0x17, 0x32, 0x4D)
SLATE = RGBColor(0x47, 0x55, 0x69)
BODY = RGBColor(0x1E, 0x29, 0x3B)
BODY2 = RGBColor(0x1B, 0x2A, 0x3A)
MUTED = RGBColor(0x8A, 0x94, 0xA3)
BORDER = RGBColor(0xD7, 0xDE, 0xE8)
ROW_ALT = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLUE = RGBColor(0x2B, 0x5F, 0x8E)
AMBER = RGBColor(0xB3, 0x6B, 0x00)
RED = RGBColor(0xB3, 0x26, 0x1E)
GREEN = RGBColor(0x2E, 0x7D, 0x32)

RAG_BG = {
    "RED": RGBColor(0xF8, 0xD8, 0xD5),
    "AMBER": RGBColor(0xFF, 0xF0, 0xD5),
    "GREEN": RGBColor(0xDD, 0xEF, 0xE2),
    "N/A": RGBColor(0xE6, 0xE6, 0xE6),
    "NO RAG": RGBColor(0xE6, 0xE6, 0xE6),
}
RAG_TEXT = {
    "RED": RED,
    "AMBER": AMBER,
    "GREEN": GREEN,
    "N/A": RGBColor(0x33, 0x33, 0x33),
    "NO RAG": BODY2,
}

FOOTER_TEXT = "Cheniere follow-on & GdF SPA economic-terms comparison"


def in_(v):
    return Emu(int(v * EMU_IN))


def add_textbox(slide, left, top, width, height, text, size, color, bold=False, font="Calibri",
                 align=PP_ALIGN.LEFT, anchor=None, wrap=True, line_spacing=None):
    box = slide.shapes.add_textbox(in_(left), in_(top), in_(width), in_(height))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if anchor is not None:
        tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return box


def add_rect(slide, left, top, width, height, fill=None, line_color=None, line_w=0.5):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, in_(left), in_(top), in_(width), in_(height))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_header(slide, title, subtitle):
    add_textbox(slide, 0.55, 0.28, 10.8, 0.42, title, 20, NAVY, bold=True, font="Aptos Display")
    add_textbox(slide, 0.58, 0.76, 11.7, 0.3, subtitle, 9.5, SLATE)
    add_rect(slide, 0.55, 1.12, 12.2, 0.0, fill=None, line_color=BORDER, line_w=1)


def add_footer(slide, page_num):
    add_textbox(slide, 0.55, 7.14, 6.8, 0.18, FOOTER_TEXT, 7.5, MUTED)
    add_textbox(slide, 12.35, 7.13, 0.4, 0.2, str(page_num), 8.5, MUTED)


def new_content_slide(prs, blank_layout, title, subtitle, page_num):
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, title, subtitle)
    add_footer(slide, page_num)
    return slide


def add_fallback_arcs(slide):
    """Approximate the Driftwood title arcs (blue large, amber mid, red small)."""
    specs = [
        (8.95, 0.62, 3.05, 3.05, BLUE, 2.2, 300),
        (9.62, 1.78, 1.55, 1.55, AMBER, 2.0, 210),
        (8.55, 2.38, 1.05, 1.05, RED, 1.8, 120),
    ]
    for left, top, w, h, color, weight, rot in specs:
        shp = slide.shapes.add_shape(MSO_SHAPE.ARC, in_(left), in_(top), in_(w), in_(h))
        shp.rotation = rot
        shp.fill.background()
        shp.line.color.rgb = color
        shp.line.width = Pt(weight)
        shp.shadow.inherit = False


def add_title_slide(prs, blank_layout, title, subtitle, tagline, purpose_label, purpose_text,
                     legend_items, page_num, arc_shapes):
    slide = prs.slides.add_slide(blank_layout)
    add_rect(slide, 0, 0, 13.333, 7.5, fill=WHITE, line_color=RGBColor(0x33, 0x33, 0x33), line_w=1)
    add_textbox(slide, 0.72, 0.78, 8.2, 0.95, title, 26, NAVY, bold=True, font="Aptos Display")
    add_textbox(slide, 0.75, 1.80, 7.9, 0.36, subtitle, 15, SLATE)
    add_textbox(slide, 0.75, 2.34, 7.4, 0.26, tagline, 10.5, BODY)
    add_textbox(slide, 0.75, 3.22, 1.3, 0.25, purpose_label, 11, BLUE, bold=True)
    add_textbox(slide, 0.75, 3.57, 6.4, 1.6, purpose_text, 13.2, BODY)
    if arc_shapes:
        for shp_xml in arc_shapes:
            slide.shapes._spTree.append(copy.deepcopy(shp_xml))
    else:
        add_fallback_arcs(slide)
    colors = [BLUE, AMBER, RED]
    xs = [8.4, 9.82, 11.55]
    widths = [1.3, 1.65, 1.3]
    for x, w, label, color in zip(xs, widths, legend_items, colors):
        add_textbox(slide, x, 5.4, w, 0.24, label, 9, color, bold=True)
    add_footer(slide, page_num)
    return slide


def add_three_cards(slide, cards, top=1.45):
    xs = [0.65, 4.85, 9.05]
    widths = [3.75, 3.75, 3.65]
    for x, w, (label, color, text) in zip(xs, widths, cards):
        add_textbox(slide, x, top, w, 0.28, label, 11.5, color, bold=True)
        add_textbox(slide, x, top + 0.36, w, 1.4, text, 10.5, BODY)


def add_rag_legend_table(slide, rows, top=3.65, row_h=0.55):
    col_x = [0.8, 2.0, 4.15]
    col_w = [1.2, 2.15, 8.2]
    headers = ["RAG", "Commercial meaning", "Key items"]
    for x, w, h in zip(col_x, col_w, headers):
        add_rect(slide, x, top, w, row_h, fill=NAVY)
        add_textbox(slide, x + 0.08, top, w - 0.16, row_h, h, 10, WHITE, bold=True, anchor=3)
    y = top + row_h
    for i, (rag, meaning, items) in enumerate(rows):
        alt = i % 2 == 1
        base_fill = ROW_ALT if alt else WHITE
        rag_bg = RAG_BG.get(rag, base_fill)
        rag_fg = RAG_TEXT.get(rag, BODY)
        add_rect(slide, col_x[0], y, col_w[0], row_h, fill=rag_bg, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[0] + 0.08, y, col_w[0] - 0.16, row_h, rag, 10, rag_fg, bold=True, anchor=3)
        add_rect(slide, col_x[1], y, col_w[1], row_h, fill=base_fill, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[1] + 0.08, y, col_w[1] - 0.16, row_h, meaning, 10, BODY, anchor=3)
        add_rect(slide, col_x[2], y, col_w[2], row_h, fill=base_fill, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[2] + 0.08, y, col_w[2] - 0.16, row_h, items, 8.5, BODY2, anchor=3)
        y += row_h


def add_two_col_table(slide, headers, rows, top=1.42, row_h=0.55, col1_w=2.45, font_size=10.2):
    col_x = [0.72, 0.72 + col1_w]
    col_w = [col1_w, 12.2 - 0.72 - col1_w - 0.15]
    for x, w, h in zip(col_x, col_w, headers):
        add_rect(slide, x, top, w, row_h, fill=NAVY)
        add_textbox(slide, x + 0.08, top, w - 0.16, row_h, h, 10, WHITE, bold=True, anchor=3)
    y = top + row_h
    for i, (label, desc) in enumerate(rows):
        alt = i % 2 == 1
        fill = ROW_ALT if alt else WHITE
        add_rect(slide, col_x[0], y, col_w[0], row_h, fill=fill, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[0] + 0.08, y, col_w[0] - 0.16, row_h, label, font_size, BODY, anchor=3)
        add_rect(slide, col_x[1], y, col_w[1], row_h, fill=fill, line_color=BORDER, line_w=0.5)
        add_textbox(slide, col_x[1] + 0.08, y, col_w[1] - 0.16, row_h, desc, font_size, BODY, anchor=3)
        y += row_h
    return y


def add_rag_table(slide, headers, rows, top=1.45, row_h=0.56, first_col_w=1.35, rag_col_w=0.9,
                   font_size=9.2, header_font_size=9.2):
    n_data_cols = len(headers) - 2
    total_w = 12.2
    remaining = total_w - first_col_w - rag_col_w - 0.15 * (len(headers) - 1)
    data_col_w = remaining / n_data_cols
    xs = [0.46]
    widths = [first_col_w] + [data_col_w] * n_data_cols + [rag_col_w]
    for w in widths[:-1]:
        xs.append(xs[-1] + w)

    for x, w, h in zip(xs, widths, headers):
        add_rect(slide, x, top, w, row_h, fill=NAVY)
        add_textbox(slide, x + 0.06, top, w - 0.12, row_h, h, header_font_size, WHITE, bold=True, anchor=3)

    y = top + row_h
    for i, row in enumerate(rows):
        alt = i % 2 == 1
        fill = ROW_ALT if alt else WHITE
        for j, (x, w) in enumerate(zip(xs, widths)):
            val = row[j]
            is_rag_col = j == len(widths) - 1
            if is_rag_col:
                rag_bg = RAG_BG.get(val, fill)
                rag_fg = RAG_TEXT.get(val, BODY)
                add_rect(slide, x, y, w, row_h, fill=rag_bg, line_color=BORDER, line_w=0.5)
                add_textbox(slide, x + 0.06, y, w - 0.12, row_h, val, font_size, rag_fg, bold=True, anchor=3)
            else:
                add_rect(slide, x, y, w, row_h, fill=fill, line_color=BORDER, line_w=0.5)
                add_textbox(slide, x + 0.06, y, w - 0.12, row_h, str(val), font_size, BODY, anchor=3)
        y += row_h
    return y


def add_badge_rows(slide, sections, top=1.35):
    y = top
    for section_title, items in sections:
        if section_title:
            add_textbox(slide, 0.75, y, 11.6, 0.24, section_title, 10.8, NAVY, bold=True)
            y += 0.36
        for rag, label, desc in items:
            bg = RAG_BG.get(rag, RGBColor(0xE6, 0xE6, 0xE6))
            fg = RAG_TEXT.get(rag, BODY)
            add_rect(slide, 0.75, y + 0.02, 0.62, 0.2, fill=bg)
            add_textbox(slide, 0.75, y + 0.02, 0.62, 0.2, rag, 7.5, fg, bold=True, align=PP_ALIGN.CENTER, anchor=3)
            add_textbox(slide, 1.49, y, 2.7, 0.25, label, 8.6, BODY, bold=True)
            add_textbox(slide, 4.25, y, 8.1, 0.28, desc, 8.4, SLATE)
            y += 0.37
    return y


def add_callout(slide, top, text, fill=RGBColor(0xEE, 0xF4, 0xF8), height=0.62):
    add_rect(slide, 0.74, top, 11.6, height, fill=fill)
    add_textbox(slide, 0.86, top, 11.4, height, text, 11.2, BODY2, anchor=3)


def add_formula_boxes(slide, boxes, top=1.25, height=2.1):
    xs = [0.55, 6.9]
    widths = [6.15, 5.9]
    for x, w, (fill, label, body) in zip(xs, widths, boxes):
        add_rect(slide, x, top, w, height, fill=fill, line_color=RGBColor(0xCB, 0xD5, 0xE1), line_w=0.75)
        add_textbox(slide, x + 0.12, top + 0.08, w - 0.24, 0.22, label, 9, BODY, bold=True)
        add_textbox(slide, x + 0.12, top + 0.36, w - 0.24, height - 0.44, body, 7.6, RGBColor(0x33, 0x41, 0x55))


def get_arc_shapes():
    if not DRIFTWOOD_PATH.exists():
        return None
    try:
        prs = Presentation(DRIFTWOOD_PATH)
    except Exception:
        return None
    slide1 = list(prs.slides)[0]
    arcs = []
    for shape in slide1.shapes:
        if shape.name in ("Shape 6", "Shape 7", "Shape 8"):
            arcs.append(shape._element)
    return arcs or None


COLS = ["Item", "BG-GC A&R", "Centrica", "GNF (CC)", "Woodside", "GdF (DES)", "RAG"]


def build():
    arc_shapes = get_arc_shapes()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    page = 1

    # ---- Slide 1: Title ----
    add_title_slide(
        prs, blank,
        "Cheniere Follow-On & GdF SPAs",
        "Economic-term comparison of the five remaining LNG contracts",
        "BG Gulf Coast A&R  •  Centrica  •  Gas Natural Fenosa (CC)  •  Woodside  •  Gaz de France",
        "Purpose",
        "Track how Cheniere's fixed-leg recovery evolved after 2011: the 2012 BG "
        "amendment scales separate UFC/MSC charges across four trains, while the "
        "2013-14 SPAs embed the fixed component as Xy and price a buyer suspension "
        "right at exactly Xy. The 2007 GdF master is the reverse trade: DES imports "
        "into Sabine Pass under seller-held cancellation options. This is an economic "
        "comparison, not a legal blackline.",
        ["Visible terms", "Structural differences", "Modelling impact"],
        page, arc_shapes,
    )
    page += 1

    # ---- Slide 2: Executive read-across ----
    slide = new_content_slide(
        prs, blank, "Executive read-across",
        "Four Cheniere FOB SPAs on the common 1.15 x HH skeleton, plus one 2007 DES master agreement; "
        "financial constants are stated in full, not redacted.",
        page,
    )
    add_three_cards(slide, [
        ("Fact", GREEN,
         "All four FOB SPAs price at 1.15 x HH. The three 2013-14 SPAs (Centrica, "
         "GNF-CC, Woodside) embed the fixed leg as Xy and give Buyer a suspension "
         "right priced at Xy x scheduled quantities. The BG A&R keeps separate "
         "UFC/MSC, stepped across Trains 1-4. Liability caps and CPI baselines "
         "are stated in full."),
        ("Inference", BLUE,
         "The 2013-14 contracts are one template family (Total-style): model "
         "suspension as a paid pause in which only the Xy leg survives. Model the "
         "BG A&R as MSC survival across cancellation, DoP and FM. Pre-5th-anniversary "
         "liability caps compute to ~US$3.29/MMBtu of ACQ in all four (arithmetic, "
         "not contract text)."),
        ("Unknown", AMBER,
         "GdF master leaves price and quantity to Specific Orders; only Specific "
         "Order No. 1 is visible. Centrica's CPI0 anchors to the year the third "
         "SPL train becomes capable: unresolvable until that year is known. "
         "Extension-period pricing is visible only in the BG A&R."),
    ])
    add_rag_legend_table(slide, [
        ("GREEN", "No material difference detected",
         "HH-linked commodity leg in all five; FOB delivery, buyer shipping, ratable scheduling in the four SPAs"),
        ("RED", "Visible economic difference",
         "UFC/MSC vs embedded Xy; suspension vs cancellation; DES vs FOB; per-train tranche mechanics (BG A&R)"),
        ("AMBER", "Framework similar, but caveats",
         "Maintenance caps and bases; CPI0 anchor conventions; credit support mechanics; GdF master openness"),
    ], top=3.55, row_h=0.5)
    page += 1

    # ---- Slide 3: Contract universe ----
    slide = new_content_slide(
        prs, blank, "Contract universe",
        "Quantity is shown but not RAG-scored; size is buyer-specific commercial volume, not a drafting difference.",
        page,
    )
    y = add_rag_table(
        slide,
        ["Contract", "Seller", "Buyer", "Signature date", "Visible quantity", "RAG"],
        [
            ["BG Gulf Coast A&R", "Sabine Pass Liquefaction", "BG Gulf Coast LNG, LLC", "25 Jan 2012",
             "286.5m MMBtu/yr at full 4-train build", "NO RAG"],
            ["Centrica", "Sabine Pass Liquefaction", "Centrica plc", "22 Mar 2013",
             "91.25m MMBtu/yr", "NO RAG"],
            ["Gas Natural Fenosa (CC)", "Corpus Christi Liquefaction", "Gas Natural Fenosa LNG SL", "2 Jun 2014",
             "78.215m MMBtu/yr", "NO RAG"],
            ["Woodside", "Corpus Christi Liquefaction", "Woodside Energy Trading Singapore", "30 Jun 2014",
             "44.12m MMBtu/yr", "NO RAG"],
            ["Gaz de France master", "Per Specific Order (SO1: GDF sells)", "Per Specific Order (SO1: Cheniere Mktg buys)",
             "26 Apr 2007", "SO1: 7 cargoes x 2.6-3.5m MMBtu", "NO RAG"],
        ],
        first_col_w=1.85, font_size=8.6, row_h=0.62,
    )
    add_callout(slide, y + 0.22,
                "Observation: the GdF contract is a DES master agreement; quantities and price sit in "
                "Specific Orders, and either party can be buyer or seller per order. The other four are "
                "20-year FOB SPAs. BG A&R ACQ builds as tranches: 182.5m (T1) + 36.5m (T2) + 34.0m (T3) "
                "+ 33.5m (T4).",
                height=0.75)
    page += 1

    # ---- Slide 4: Three structural families ----
    slide = new_content_slide(
        prs, blank, "Three structural families",
        "The five contracts split into a multi-train UFC/MSC amendment, an embedded-Xy suspension template, "
        "and a DES master framework.",
        page,
    )
    add_three_cards(slide, [
        ("1. Multi-train UFC / MSC", BLUE,
         "BG Gulf Coast A&R (2012). Separate capacity charge: UFC Base steps up at "
         "each Train DFCD (1.9125 / 2.0375 / 2.1215 / 2.1848); MSC = UFC x Q/12 "
         "invoiced monthly; per-train tranche cancellation rights."),
        ("2. Embedded Xy + suspension", AMBER,
         "Centrica (2013), GNF-CC and Woodside (2014). Fixed leg inside CSP as Xy; "
         "buyer may suspend deliveries (min. 1 month, 2-month notice) paying a "
         "Suspension Fee = Xy x scheduled quantities of the suspended month."),
        ("3. DES master + Specific Orders", RED,
         "GdF (2007). Framework contract; each Specific Order sets parties, "
         "quantity, price and ships. SO1: GDF sells 7 cargoes DES into Sabine Pass "
         "at 94% x HH minus US$0.65, with seller-held cancellation options."),
    ], top=1.5)
    add_callout(slide, 3.5,
                "Lineage: the 2011 SPL SPAs recover the fixed leg as UFC/MSC (BG, GNF) or embedded Xy "
                "(GAIL, Total). The 2012 BG A&R extends UFC/MSC to four trains. From Centrica (2013) "
                "onward the embedded-Xy form wins, and the Total-style suspension right becomes "
                "standard; Corpus Christi raises X0 from US$3.00 to US$3.50.",
                height=0.95)
    add_callout(slide, 4.7,
                "Modelling consequence: the same economic object (a CPI-escalated fixed leg that survives "
                "buyer non-lifting) is invoiced through three different routes: separate monthly charge, "
                "cargo price component, or suspension fee.", fill=RGBColor(0xF5, 0xF7, 0xFA), height=0.8)
    page += 1

    # ---- Slide 5: Visible price formulas ----
    slide = new_content_slide(
        prs, blank, "Visible price formulas",
        "All five are Henry Hub-linked and stated in full; architecture, constants and CPI weights differ.",
        page,
    )
    y = add_rag_table(
        slide,
        ["Contract", "Visible formula", "Fixed / capacity component", "RAG"],
        [
            ["BG Gulf Coast A&R", "CSP = 1.15 x HH",
             "UFCm = Base + 0.3375 x CPIy/CPI0; Base steps 1.9125 / 2.0375 / 2.1215 / 2.1848 by Train DFCD; MSC = UFC x Q/12", "RED"],
            ["Centrica", "CSP = (1.15 x HH) + Xy",
             "Xy = (0.885 + 0.115 x CPI(y-1)/CPI0) x X0; X0 = US$3.00/MMBtu", "RED"],
            ["GNF (Corpus Christi)", "CSP = (1.15 x HH) + Xy",
             "Xy = (0.86 + 0.14 x CPI(y-1)/CPI0) x X0; X0 = US$3.50/MMBtu", "RED"],
            ["Woodside", "CSP = (1.15 x HH) + Xy",
             "Xy = (0.885 + 0.115 x CPI(y-1)/CPI0) x X0; X0 = US$3.50/MMBtu", "RED"],
            ["GdF Specific Order 1", "P = 94% x HH (NYMEX prompt) - US$0.65",
             "None; flat per-order price (master: price per Specific Order)", "RED"],
        ],
        first_col_w=1.85, rag_col_w=0.85, font_size=8.6, row_h=0.68,
    )
    add_callout(slide, y + 0.2,
                "RAG is RED because architecture differs, not the HH link: the four SPAs sell at a "
                "premium to HH (export economics); GdF SO1 buys at a discount to HH (2007 import "
                "economics into the same terminal). No price term on this slide is redacted.",
                height=0.75)
    page += 1

    # ---- Slide 6: Appendix: formula detail as filed ----
    slide = new_content_slide(
        prs, blank, "Appendix A: formula detail as filed",
        "CPI series is BLS CPI-U (CUUR0000SA0) throughout; the CPI0 anchor convention differs per contract.",
        page,
    )
    add_formula_boxes(slide, [
        (RGBColor(0xDB, 0xEA, 0xFE), "BG Gulf Coast A&R: UFC / MSC (9.2, 9.3)",
         "UFCm = Base + (0.3375 x CPIy / CPI0)\n"
         "Base: 1.9125 pre-T2 DFCD; 2.0375 T2; 2.1215 T3; 2.1848 T4\n"
         "MSC = UFC x Q/12, Q = MACQ - MQ (MACQ steps by train)\n"
         "Extension Base: 1.9125 up to 182.5m MMBtu, above that\n"
         "Base = (182.5m/A) x 1.9125 + ((A-182.5m)/A) x 2.6625\n"
         "CPI0 = 12-month average preceding Train 1 DFCD"),
        (RGBColor(0xFE, 0xE2, 0xE2), "Embedded Xy family (9.1.2)",
         "Centrica:  Xy = (0.885 + 0.115 x CPI(y-1)/CPI0) x US$3.00\n"
         "GNF-CC:   Xy = (0.86 + 0.14 x CPI(y-1)/CPI0) x US$3.50\n"
         "Woodside: Xy = (0.885 + 0.115 x CPI(y-1)/CPI0) x US$3.50\n\n"
         "CPI0 anchors: Centrica = calendar year in which the third\n"
         "SPL train becomes capable; GNF-CC = calendar 2017;\n"
         "Woodside = calendar 2014"),
    ], top=1.3, height=2.15)
    add_formula_boxes(slide, [
        (RGBColor(0xEE, 0xF4, 0xF8), "GdF Specific Order 1: price and underdelivery",
         "P = 94% x NYMEX HH final settlement (Prompt Month) - US$0.65\n"
         "Prompt Month = month in which the Arrival Period ends\n"
         "(next month if it ends on/after the 25th)\n\n"
         "Cargo Underdelivery Amount = Nominal Quantity x\n"
         "max(Cover Difference, US$1.00); Cover Difference =\n"
         "4-day Platts Gas Daily HH average after the Arrival Period\n"
         "minus 80% x P"),
        (RGBColor(0xF5, 0xF7, 0xFA), "GdF Specific Order 1: cancellation fees",
         "Lot 1: either party may cancel by 29 Feb 2008, no liability.\n"
         "Seller may cancel lots 2-7:\n"
         "  all six: 0.25 USD x 6 x 3.05m MMBtu (Total Cancellation Fee)\n"
         "  any one: 0.50 USD x 3.05m MMBtu (Individual Cancellation Fee)\n"
         "(3.05m = (2.6m + 3.5m)/2 deemed Nominal Quantity)\n"
         "Fees are Buyer's sole and exclusive remedy; no replacement\n"
         "cargo obligation"),
    ], top=3.7, height=2.15)
    add_callout(slide, 6.0,
                "Fact: additive CPI escalation (BG A&R: 0.3375 term on top of a fixed Base) and "
                "multiplicative escalation (Xy family: 11.5-14% weight inside the constant) respond "
                "differently to inflation; do not treat the fixed legs as interchangeable.",
                height=0.62)
    page += 1

    # ---- Slide 7: Fixed-leg survival on non-lifting ----
    slide = new_content_slide(
        prs, blank, "Non-lifting rights and fixed-leg survival",
        "In all four Cheniere SPAs the fixed leg survives buyer non-lifting; only the invoicing route differs.",
        page,
    )
    y = add_rag_table(slide, COLS, [
        ["Buyer non-lifting right",
         "Cancellation Right (5.6): no fee, any cargo, notice by Cancellation Deadline",
         "Suspension (5.7): whole months, 2-month notice, min. 1 month",
         "Suspension (5.7): same template", "Suspension (5.7): same template",
         "None (SO1 gives Seller the cancellation options)", "RED"],
        ["Cost of not lifting",
         "MSC continues in full",
         "Suspension Fee = Xy x sum of SCQ in the suspended month",
         "Suspension Fee = Xy x sum of SCQ", "Suspension Fee = Xy x sum of SCQ",
         "n/a for Buyer; Seller pays fee if it cancels", "RED"],
        ["Fixed leg under seller failure (DoP)",
         "DoP priced off CSP minus UFC; MSC still payable",
         "Cargo DoP priced off CSP (Xy inside)",
         "Cargo DoP priced off CSP", "Cargo DoP priced off CSP",
         "Underdelivery Amount = NQ x max(HH cover diff, US$1.00)", "RED"],
        ["Fixed leg under FM",
         "MSC reduced prospectively in proportion to FM quantities",
         "Not extracted", "Not extracted", "Not extracted", "Not extracted", "AMBER"],
    ], first_col_w=1.7, rag_col_w=0.75, font_size=7.5, row_h=0.86)
    add_callout(slide, y + 0.18,
                "Do not value suspension as free volume flexibility: the Suspension Fee is exactly the "
                "fixed leg on the suspended quantities, so suspension monetises only the commodity leg. "
                "Economically it matches the BG cancellation + surviving MSC.", height=0.72)
    page += 1

    # ---- Slide 8: Term and extension mechanics ----
    slide = new_content_slide(
        prs, blank, "Term and extension mechanics",
        "The four SPAs share the 20-year + up-to-10-year extension skeleton; the BG A&R adds per-train "
        "tranche cancellation; the GdF master is open-ended.",
        page,
    )
    y = add_rag_table(slide, COLS, [
        ["Base term",
         "20 years from Train 1 DFCD", "20 years from DFCD", "20 years from DFCD", "20 years from DFCD",
         "Through the GDF Transatlantic Option Agreement period, then evergreen", "AMBER"],
        ["Extension",
         "Up to 10 years, any portion of ACQ; per-tranche reductions allowed",
         "Up to 10 years, any portion of ACQ; notice by 17th anniversary of DFCD",
         "Same as Centrica", "Same as Centrica",
         "30-day notice termination after the primary term", "AMBER"],
        ["Extension-period price",
         "UFC Base 1.9125 up to 182.5m MMBtu; blended to 2.6625 above",
         "Not extracted", "Not extracted", "Not extracted", "n/a", "AMBER"],
        ["Project-linkage optionality",
         "Trains 2-4: FID/DFCD-linked tranche cancellation rights (4.6), either party",
         "Single-train linkage; DFCD windows", "Single-train linkage", "Single-train linkage",
         "SO1 lot 1: mutual cancellation right to 29 Feb 2008", "RED"],
    ], first_col_w=1.7, rag_col_w=0.75, font_size=7.5, row_h=0.86)
    add_callout(slide, y + 0.18,
                "The BG A&R embeds a strip of train-by-train real options on both sides: each incremental "
                "tranche (36.5m / 34.0m / 33.5m MMBtu) can fall away if the corresponding FID or DFCD "
                "fails, with UFC Base and MACQ stepping accordingly.", height=0.72)
    page += 1

    # ---- Slide 9: Volume, shape and maintenance ----
    slide = new_content_slide(
        prs, blank, "Volume, shape and maintenance",
        "Maintenance moves from percentage caps at Sabine Pass to the 7.5% / 25%-in-6-years form at "
        "Corpus Christi (the form Driftwood later uses).",
        page,
    )
    y = add_rag_table(slide, COLS, [
        ["Annual quantity",
         "ACQ = sum of train tranches; 286.5m MMBtu at full build",
         "91.25m MMBtu", "78.215m MMBtu", "44.12m MMBtu",
         "Per Specific Order (SO1: 7 x 2.6-3.5m MMBtu)", "NO RAG"],
        ["Delivery shape",
         "Even and ratable ADP", "Even and ratable ADP", "Even and ratable ADP", "Even and ratable ADP",
         "SO1: one cargo per month, Apr-Oct 2008", "GREEN"],
        ["Major maintenance",
         "6%/yr of then-applicable ACQ; cumulative caps step 54.75m to 83.25m MMBtu by train stage",
         "6%/yr; 30% of ACQ over the initial Term",
         "7.5%/yr; 25% over any 6 consecutive Contract Years",
         "7.5%/yr; 25% over any 6 consecutive Contract Years",
         "n/a", "AMBER"],
        ["Inspection allowance",
         "Per-train absolute caps: 15.15m / 4.56m / 2.83m / 5.58m MMBtu in defined windows",
         "8.3%/yr; 17% of ACQ over the initial Term",
         "None separate from major maintenance", "None separate", "n/a", "AMBER"],
    ], first_col_w=1.6, rag_col_w=0.75, font_size=7.4, row_h=0.9)
    add_callout(slide, y + 0.15,
                "Centrica's 17% cumulative inspection cap matches Total (2012), not the 16.6% in the "
                "2011 BG/GNF SPAs. The Corpus Christi 7.5% / 25%-in-6-years form reappears in the 2021 "
                "Driftwood SPAs.", height=0.62)
    page += 1

    # ---- Slide 10: Damages, caps and credit ----
    slide = new_content_slide(
        prs, blank, "Damages, liability caps and credit",
        "Caps are stated in full; pre-5th-anniversary caps compute to ~US$3.29 per MMBtu of ACQ in all "
        "four SPAs (computed ratio, not contract text).",
        page,
    )
    y = add_rag_table(slide, COLS, [
        ["Seller liability cap (pre-5th anniv.)",
         "USD 600m / 720m / 832m / 942m by Train FID stage",
         "USD 300m", "USD 257.145m", "USD 145.052m", "Not extracted", "AMBER"],
        ["Seller liability cap (after)",
         "2 x ACQ figure in MMBtu (e.g. 573m at full build)",
         "USD 400m", "USD 342.86m", "USD 193.402m", "Not extracted", "AMBER"],
        ["Cover damages (buyer failure)",
         "CSP-based; cancelled-late quantities become Cargo Shortfall; MSC unaffected",
         "Cover Damages, CSP-based", "Same template", "Same template",
         "n/a in same form", "AMBER"],
        ["Buyer credit support",
         "Buyer-specific mechanics; not extracted in detail",
         "Parent guaranty (Exhibit C) on merger/asset sale; Centrica plc itself is Buyer",
         "15.3 credit support clause; details not extracted",
         "15.3 credit support clause; details not extracted",
         "Assignment to banks >= US$1bn assets allowed for financing", "AMBER"],
    ], first_col_w=1.85, rag_col_w=0.75, font_size=7.4, row_h=0.9)
    add_callout(slide, y + 0.15,
                "Computed ratios: 300/91.25 = 257.145/78.215 = 145.052/44.12 = 3.2877 USD/MMBtu; "
                "post-anniversary caps = 4.3836 USD/MMBtu for the 2013-14 family but 2.00 USD/MMBtu "
                "in the BG A&R. The cap schedule is formulaic across the portfolio.", height=0.72)
    page += 1

    # ---- Slide 11: GdF DES mechanics ----
    slide = new_content_slide(
        prs, blank, "GdF master: DES mechanics",
        "The only DES contract in the collection; freight, laytime and demurrage sit with the seller side, "
        "reversing the FOB allocation of the SPAs.",
        page,
    )
    y = add_two_col_table(slide, ["Item", "As filed"], [
        ("Structure", "Master framework + Specific Orders; either party may be buyer or seller per order"),
        ("Title and risk", "Pass at the Delivery Point on unloading (DES); vapour return passes buyer to seller"),
        ("SO1 trade", "GDF International Trading sells 7 cargoes to Cheniere Marketing, DES Sabine Pass Terminal, Apr-Oct 2008"),
        ("Laytime / demurrage", "36 hours up to 145,000 m3 (formula above); demurrage US$100,000 per 24h, prorated"),
        ("Payment", "20 Business Days from invoice; disputed invoices paid 100% provisionally"),
        ("Late payment interest", "2% over Euribor (EUR) or Libor (USD), 365-day year"),
        ("Nominations", "Seller nominates quantity, ship, heel, loading port and arrival day by the 1st of the prior month"),
        ("Governing law (letter agt.)", "England and Wales; expert determination via ICC International Centre for Expertise"),
    ], row_h=0.52, font_size=9.2, col1_w=2.6)
    add_callout(slide, y + 0.15,
                "Same physical terminal as the SPL SPAs, reversed flow: a 2007 import trade priced below "
                "HH (94% x HH - 0.65) versus post-2011 export trades priced above HH (1.15 x HH + fixed "
                "leg). Useful as a pre-shale-reversal benchmark.", height=0.72)
    page += 1

    # ---- Slide 12: First-pass RAG summary ----
    slide = new_content_slide(
        prs, blank, "First-pass RAG summary",
        "High-confidence differences concentrate in price architecture, non-lifting mechanics, "
        "train linkage and delivery basis.",
        page,
    )
    add_rag_table(slide, COLS, [
        ["Signature date", "25 Jan 2012", "22 Mar 2013", "2 Jun 2014", "30 Jun 2014", "26 Apr 2007", "NO RAG"],
        ["Commodity index", "HH", "HH", "HH", "HH", "HH (NYMEX prompt)", "GREEN"],
        ["Price architecture", "HH + UFC/MSC", "HH + embedded Xy", "HH + embedded Xy", "HH + embedded Xy",
         "Flat per-order (94% HH - 0.65)", "RED"],
        ["Fixed-leg constant", "UFC Base 1.9125-2.1848", "X0 = 3.00", "X0 = 3.50", "X0 = 3.50", "None", "RED"],
        ["CPI escalation", "Additive 0.3375 term", "11.5% weight", "14% weight", "11.5% weight", "None", "RED"],
        ["Non-lifting mechanics", "Cancellation + MSC", "Suspension fee", "Suspension fee", "Suspension fee",
         "Seller cancellation fees", "RED"],
        ["Delivery basis", "FOB Sabine Pass", "FOB Sabine Pass", "FOB Corpus Christi", "FOB Corpus Christi",
         "DES Sabine Pass (SO1)", "RED"],
        ["Maintenance form", "6% + train caps", "6%/30%; 8.3%/17%", "7.5%/25% in 6y", "7.5%/25% in 6y", "n/a", "AMBER"],
    ], first_col_w=1.6, rag_col_w=0.75, font_size=7.4, row_h=0.6)
    page += 1

    # ---- Slide 13: Recommended valuation object model ----
    slide = new_content_slide(
        prs, blank, "Recommended valuation object model",
        "Separate cargo economics from capacity economics before computing value, hedge and risk.",
        page,
    )
    add_three_cards(slide, [
        ("1. Variable commodity leg", BLUE,
         "1.15 x HH x loaded MMBtu in the four SPAs; 0.94 x HH - 0.65 in GdF SO1 "
         "(buy side)."),
        ("2. Fixed liquefaction leg", AMBER,
         "UFC/MSC with train-stepped Base (BG A&R) or CPI-escalated Xy "
         "(3.00 / 3.50 base). None in GdF."),
        ("3. Non-lifting logic", RED,
         "MSC survival; Suspension Fee = Xy x SCQ; per-train tranche cancellation; "
         "GdF seller cancellation fees."),
    ], top=1.5)
    add_callout(slide, 3.15,
                "Inputs to store per contract: train tranches and DFCD triggers; ACQ; UFC Base schedule "
                "or X0 and CPI weight; CPI0 anchor; suspension fee formula; cancellation deadlines and "
                "fees; DoP reference price (CSP vs CSP minus UFC); liability cap schedule; delivery "
                "basis (FOB/DES).", fill=RGBColor(0xF5, 0xF7, 0xFA), height=1.05)
    add_callout(slide, 4.4,
                "Output metrics: margin by leg; fixed-leg recovery under non-lifting scenarios; hedge "
                "exposure to HH and CPI; train-linkage option values (BG A&R); cap-constrained DoP "
                "exposure; credit-adjusted PV.", fill=RGBColor(0xEE, 0xF4, 0xF8), height=0.85)
    page += 1

    # ---- Slide 14: Caveats and next steps ----
    slide = new_content_slide(
        prs, blank, "Caveats and next steps",
        "This is an economic comparison; it is not a legal blackline.",
        page,
    )
    add_badge_rows(slide, [
        (None, [
            ("AMBER", "GdF master openness", "Only Specific Order No. 1 is visible; other orders may differ in "
             "price, quantity and direction."),
            ("AMBER", "Centrica CPI0 anchor", "Anchored to the year the third SPL train becomes capable; value of "
             "Xy is unresolved until that year is fixed."),
            ("AMBER", "BG A&R MACQ / tranche interaction", "MSC quantity Q = MACQ - MQ steps with train DFCDs and "
             "tranche cancellations; model clause-by-clause."),
            ("AMBER", "Extension-period pricing", "Visible only in the BG A&R (Base blend to 2.6625); not extracted "
             "for the 2013-14 family."),
            ("AMBER", "FM effect on fixed leg", "Extracted only for the BG A&R (MSC reduced prospectively); check "
             "14.x in the other three SPAs."),
            ("AMBER", "Credit support thresholds", "Parent-guaranty forms identified; ratings triggers and LC "
             "mechanics not extracted."),
        ]),
    ], top=1.5)
    add_callout(slide, 4.35,
                "Main structural addition to the master comparison: record the invoicing route of the "
                "fixed leg (separate MSC vs embedded Xy vs suspension fee) before valuing non-lifting.",
                fill=RGBColor(0xF5, 0xF7, 0xFA), height=0.6)
    page += 1

    OUT_PATH.parent.mkdir(exist_ok=True)
    prs.save(OUT_PATH)
    print(f"saved {OUT_PATH} ({page - 1} slides)")


if __name__ == "__main__":
    build()
